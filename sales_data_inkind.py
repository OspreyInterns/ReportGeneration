
import sqlite3 as sqlite
import logging
import openpyxl
from openpyxl.styles import Alignment
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

# case column numbers
CMSW_CASE_ID = 0
SERIAL_NUMBER = 3
DATE_OF_PROCEDURE = 5
DYEVERT_USED = 6
THRESHOLD_VOLUME = 8
ATTEMPTED_CONTRAST_INJECTION_VOLUME = 13
DIVERTED_CONTRAST_VOLUME = 14
CUMULATIVE_VOLUME_TO_PATIENT = 15
PERCENTAGE_CONTRAST_DIVERTED = 16
# Total duration.  
# 19 for CMSW 
# 20 for iPad 
TOTAL_DURATION = 20


# colors
WHITE = 0
LIGHT_GREEEN = 1
GREEN = 2
YELLOW = 3
RED = 4



# Write data for sales team to appropriate templates for power point and excel


def _sort_criteria(case):
    """reads info for a sort"""
    return case[1], case[2]


def list_builder(file_names):
    """Takes the list of files and builds the list of lists to write"""
    print('building list')
    cases = []
    inkind=0
    divertedvolcount =0

    for file_name in file_names:
        con = sqlite.connect(file_name)

        with con:

            cur = con.cursor()
            cur.execute('SELECT * FROM CMSWCases')
            rows = cur.fetchall()

            for row in rows:
                
                # ~ if row[DYEVERT_USED] == 1:
                    # if row[THRESHOLD_VOLUME] == 0:
                        # debug_msg = 'CMSW ' + str(row[SERIAL_NUMBER]) + ', case ' +\
                                    # str(row[CMSW_CASE_ID]) + ' has zero threshold'
                        # logging.warning(debug_msg)
                        # print(debug_msg)
                # adjust the in-kind logic here.  Create list cases_excluded with reason as last column

                # if not(row[TOTAL_DURATION] <= 5 or row[ATTEMPTED_CONTRAST_INJECTION_VOLUME] ==
                # row[DIVERTED_CONTRAST_VOLUME] == row[CUMULATIVE_VOLUME_TO_PATIENT] ==
                # row[PERCENTAGE_CONTRAST_DIVERTED] == 0):

                # use this block for traffic-light over/under thresholdred color scheme

                # if row[CUMULATIVE_VOLUME_TO_PATIENT] <= row[THRESHOLD_VOLUME] :
                    # color = GREEN
                # # elif row[CUMULATIVE_VOLUME_TO_PATIENT] <= row[THRESHOLD_VOLUME]:
                    # # color = GREEN
                # elif row[CUMULATIVE_VOLUME_TO_PATIENT] <= row[THRESHOLD_VOLUME]*3/1.5:
                    # color = YELLOW
                # elif row[CUMULATIVE_VOLUME_TO_PATIENT] >= row[THRESHOLD_VOLUME] *3/1.5 :
                    # color = RED

                # else:
                # color = WHITE

# use this block for 5 color default scheme:

                if row[CUMULATIVE_VOLUME_TO_PATIENT] <= row[THRESHOLD_VOLUME] \
                        / 3 <= row[ATTEMPTED_CONTRAST_INJECTION_VOLUME]:
                    color = LIGHT_GREEEN
                elif row[CUMULATIVE_VOLUME_TO_PATIENT] <= row[THRESHOLD_VOLUME] \
                        * 2 / 3 <= row[ATTEMPTED_CONTRAST_INJECTION_VOLUME]:
                    color = GREEN
                elif row[CUMULATIVE_VOLUME_TO_PATIENT] <= row[THRESHOLD_VOLUME] \
                        <= row[ATTEMPTED_CONTRAST_INJECTION_VOLUME]:
                    color = YELLOW
                elif row[CUMULATIVE_VOLUME_TO_PATIENT] >= row[THRESHOLD_VOLUME] \
                        <= row[ATTEMPTED_CONTRAST_INJECTION_VOLUME]:
                    color = RED

                else:
                    color = WHITE

                #in kind comments criteria
                comment=''
                if row[TOTAL_DURATION] <=5 :
                    comment='Case less than 5 minutes'
                elif row[CUMULATIVE_VOLUME_TO_PATIENT] == row[DIVERTED_CONTRAST_VOLUME] == row[ATTEMPTED_CONTRAST_INJECTION_VOLUME]==0:
                    comment='0 mL contrast injected'
                elif row[DYEVERT_USED]== 0 :
                    comment='DyeTect Case'
                elif row[DIVERTED_CONTRAST_VOLUME] < 5 :
                    comment ='Diverted Volume < 5 mL'
                elif row[ATTEMPTED_CONTRAST_INJECTION_VOLUME] < 20 :
                    comment ='Attempted Volume < 20 mL'
                elif row[CUMULATIVE_VOLUME_TO_PATIENT] < 30 :
# No In-Kind anymore
#                    comment = 'In-Kind for Kidneys 2.0 case'
#                    inkind =inkind +1
# check for clinical diverted volume <10%
                    if row[PERCENTAGE_CONTRAST_DIVERTED] < 10:
                        divertedvolcount = divertedvolcount + 1
                        comment += ", Diverted volume < 10%"
                else:
                    comment =''
                        
# add logic here for in-kind count.  if meets IK2.0 then reason as last string item in the list.  Mod append for new string. 
                cases.append((color, row[DATE_OF_PROCEDURE][0:10], row[DATE_OF_PROCEDURE][11:22],
                             row[THRESHOLD_VOLUME], row[ATTEMPTED_CONTRAST_INJECTION_VOLUME],
                             row[CUMULATIVE_VOLUME_TO_PATIENT], row[DIVERTED_CONTRAST_VOLUME],
                             row[PERCENTAGE_CONTRAST_DIVERTED],comment))
                                    

                    
                                        
    cases.sort(key=_sort_criteria)
    return cases


def write(file_names, cmsw):
    """Writes data into the an Excel Sheet and a Power Point slide as seen in the example
    Takes two inputs:
        -the file names of the CMSW databases
        -the serial numbers of the CMSWs
    Generates two files:
        -The summary table
        -The Power Point slide, which was copied from the example
    the Excel In-Kind template is different from the Sale report template:
        - case data begins on line 20.    
        - In-Kind summary table 
        - Diverted Vol <10% count for Clinical team. 
    """
    
    # this will need adjustment to write the cases, then the cases_excluded list, with the summary table. Check if
    # insertions work on the excel library.
    
    print('Processing data for sales report')
    cases = list_builder(file_names)
    xlsx_name = str(cmsw) + '-data-tables-inkind.xlsx'
    wb = openpyxl.load_workbook('Sales-Template-inkind.xlsx')
    data_sheet = wb.active
    data_sheet.title = 'Sheet1'
    print('Data ready, writing sales report')


    for row in range(len(cases)):
        for col in range(len(cases[row])):
            data_sheet.cell(row=row + 20, column=col + 1, value=cases[row][col])
            data_sheet.cell(row=row + 20, column=col + 1).alignment = Alignment(wrapText=True)

    data_sheet.column_dimensions['A'].hidden = False
    wb.save(xlsx_name)

    pptx_name = str(cmsw) + '-slide.pptx'
    colors = [0, 0, 0, 0]
    attempted = 0
    diverted = 0
    print('Report written, constructing slide')
    for case in cases:
        if case[0] == RED:
            colors[0] += 1
        elif case[0] == YELLOW:
            colors[1] += 1
        elif case[0] == GREEN:
            colors[2] += 1
        elif case[0] == LIGHT_GREEEN:
            colors[3] += 1
        attempted += case[4]
        diverted += case[6]

    percent_saved = round(diverted/attempted*100)
    prs = Presentation('Slide-Template.pptx')
    total = colors[0] + colors[1] + colors[2] + colors[3]
    title = 'N=' + str(total)
    data = CategoryChartData()
    data.add_series(title, colors)
    data.categories = ['> Threshold, N=', '< Threshold, N=', '< 2/3 Threshold, N=', '< 1/3 Threshold, N=']
    prs.slides[0].shapes[4].chart.replace_data(data)
    text = ['All cases (N=' + str(len(cases)) + ')',
            str(percent_saved) + '% avg \nLess \nContrast',
            str(round(diverted)) + ' mL less total']

    for box in range(5, 8, 1):
        prs.slides[0].shapes[box].text_frame.clear()
        prs.slides[0].shapes[box].text_frame.paragraphs[0].text = text[box-5]
        prs.slides[0].shapes[box].text_frame.paragraphs[0].allignment = PP_ALIGN.CENTER
        prs.slides[0].shapes[box].text_frame.paragraphs[0].font.size = Pt(26-2*box)

    prs.slides[0].shapes[5].text_frame.paragraphs[0].font.bold = True
    prs.slides[0].shapes[7].text_frame.paragraphs[0].font.italic = True
    prs.save(pptx_name)
    print('Slides complete')
